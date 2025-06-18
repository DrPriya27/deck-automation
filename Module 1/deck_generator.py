###############################################################################  
# ppt_generator.py  –  AI slide-deck generator that outputs PowerPoint (.pptx)  
###############################################################################  
import os, base64, mimetypes, json, textwrap, re, io, argparse
from pathlib import Path  
  
from dotenv import load_dotenv  
from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation                 # pip install python-pptx  
from pptx.util import Inches, Pt  
  
###############################################################################  
# 1.  Azure OpenAI client  
###############################################################################  
load_dotenv()  

# Initialize Azure OpenAI client
llm = AzureChatOpenAI(
            azure_endpoint=os.getenv('AZURE_OPENAI_ENDPOINT'),
            azure_deployment="gpt-4o",
            api_version="2024-06-01",
            api_key=os.getenv('AZURE_OPENAI_API_KEY'),
            max_tokens=None,
            temperature=0.7,
)
  
###############################################################################  
# 2.  Helper functions  
###############################################################################  
IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}  
  
def image_to_b64(p: Path) -> str:  
    return base64.b64encode(p.read_bytes()).decode("utf-8")  
  
def read_folder(folder: Path):  
    """Return (markdown_text, images_dict {filename: b64}) for every file."""  
    md_chunks, image_dict = [], {}  
    for fp in folder.rglob("*"):  
        if fp.suffix.lower() == ".md":  
            md_chunks.append(fp.read_text(encoding="utf-8", errors="ignore"))  
        elif fp.suffix.lower() in IMG_EXT:  
            image_dict[fp.name] = image_to_b64(fp)  
    return "\n\n".join(md_chunks), image_dict  
  
def build_initial_messages(prompt: str, md: str, images: dict):  
    """First call – supply markdown text + pictures as context."""  
    sys_msg = {  
        "role": "system",  
        "content": textwrap.dedent("""  
            You are SlideBuilder-GPT.  
            Return ONLY valid JSON (no markdown, no code fences).  
            JSON schema:  
              {  
                "slides":[  
                  {  
                    "title":"Slide title",  
                    "points":["bullet 1","bullet 2"],  
                    "images":["file1.png","file2.jpg"],        // optional  
                    "notes":"optional speaker notes"  
                  }  
                ]  
              }  
            Rules:  
            • The file names in the images array MUST match exactly the files  
              you have been shown (case-sensitive).  
            • Do NOT embed base-64.  Do NOT output HTML.  
        """).strip(),  
    }  
  
    user_parts = [{"type": "text", "text": prompt}]  
    if md.strip():  
        user_parts.append({"type": "text",  
                           "text": f"\n\nHere are the markdown files:\n\n{md}"})  
    for fn, b64 in images.items():  
        mime, _ = mimetypes.guess_type(fn)  
        mime = mime or "image/png"  
        # show model the actual picture  
        user_parts.append({  
            "type": "image_url",  
            "image_url": {"url": f"data:{mime};base64,{b64}"},  
        })  
        user_parts.append({"type": "text", "text": f"(filename: {fn})"})  
  
    return [sys_msg, {"role": "user", "content": user_parts}]  
  
def openai_call(messages):  
    messages_mod = [
            SystemMessage(
                content="You are a slide deck generator that creates PowerPoint presentations. Return only valid JSON following the specified format."
            ),
            HumanMessage(content=messages),
    ]

    resp = llm.invoke(messages)
    
    return resp.content.strip()  
  
def parse_slides_json(raw: str) -> dict:  
    """  
    Ensure the assistant reply is valid JSON.  
    If the model wraps it in ```json``` fences, strip them.  
    """  
    raw = raw.strip()  
    if raw.startswith("```"):  
        raw = re.sub(r"^```[a-zA-Z]*", "", raw).rstrip("`").strip()  
    return json.loads(raw)  
  
# -------------------------- PowerPoint generation -------------------------  
  
def build_pptx(slides_dict: dict, folder: Path) -> bytes:  
    prs = Presentation()  
    title_content_layout = prs.slide_layouts[1]   # Title and Content  
  
    for slide_data in slides_dict.get("slides", []):  
        title_text   = slide_data.get("title", "")  
        bullets      = slide_data.get("points", []) or []  
        images       = slide_data.get("images", []) or []  
        notes_text   = slide_data.get("notes", "")  
  
        slide = prs.slides.add_slide(title_content_layout)  
  
        # Title  
        slide.shapes.title.text = title_text  
  
        # Bullets  
        body = slide.shapes.placeholders[1].text_frame  
        body.clear()                 # remove default bullet  
        for idx, b in enumerate(bullets):  
            p = body.add_paragraph() if idx else body.paragraphs[0]  
            p.text  = b  
            p.level = 0  
            p.font.size = Pt(18)  
  
        # Images – stack under the body placeholder  
        pic_left = Inches(5.5)      # right-hand side  
        pic_top  = Inches(1.5)  
        maxw     = Inches(3.0)  
        for img in images:  
            img_path = folder / img  
            if img_path.exists():  
                slide.shapes.add_picture(str(img_path), pic_left, pic_top,  
                                          width=maxw)  
                pic_top += Inches(2.5)  
  
        # Speaker notes  
        if notes_text:  
            notes = slide.notes_slide.notes_text_frame  
            notes.text = notes_text  
  
    return prs

def save_presentation(prs, output_path):
    """Save the presentation to the specified path"""
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def main():
    # Set up argument parser
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentations using AI')
    parser.add_argument('--folder', '-f', required=True, help='Path to folder with markdown and images')
    parser.add_argument('--prompt', '-p', required=True, help='Prompt describing what slides to generate')
    parser.add_argument('--output', '-o', default='deck.pptx', help='Output PowerPoint file name')
    args = parser.parse_args()
    
    # Convert folder path to Path object and validate
    folder_path = Path(args.folder).expanduser()
    if not folder_path.is_dir():
        print(f"Error: {args.folder} is not a valid directory")
        return
    
    # Read content from the folder
    print(f"Reading content from {folder_path}...")
    md_text, images_dict = read_folder(folder_path)
    
    # Build the initial messages
    print("Building request for AI model...")
    openai_messages = build_initial_messages(args.prompt, md_text, images_dict)
    
    # Call OpenAI
    print("Generating slide content with AI...")
    raw_reply = openai_call(openai_messages)
    
    # Parse the JSON response
    try:
        print("Parsing AI response...")
        slides_dict = parse_slides_json(raw_reply)
        
        # Build the PowerPoint presentation
        print("Building PowerPoint presentation...")
        prs = build_pptx(slides_dict, folder_path)
        
        # Save the presentation
        output_path = args.output
        save_presentation(prs, output_path)
        
    except Exception as ex:
        print(f"Error: Could not build PowerPoint: {ex}")
        return

if __name__ == "__main__":
    main()
