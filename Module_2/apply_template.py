"""
Simple script to apply a PowerPoint template to an existing presentation.
"""
import os
import sys
from pptx import Presentation

def copy_slides_to_template(source_path, template_path, output_path):
    """
    Copy slides from source presentation to template.
    """
    print(f"Checking source file: {source_path}")
    if not os.path.exists(source_path):
        print(f"ERROR: Source file not found: {source_path}")
        return False
        
    print(f"Checking template file: {template_path}")
    if not os.path.exists(template_path):
        print(f"ERROR: Template file not found: {template_path}")
        return False
        
    try:
        # Load presentations
        print("Loading source presentation...")
        source = Presentation(source_path)
        
        print("Loading template presentation...")
        template = Presentation(template_path)
        
        # Save template slides for reference
        print("Creating new presentation...")
        template_slides = [slide for slide in template.slides]
        
        # Create output presentation from template
        output = Presentation(template_path)
        
        # Remove existing slides from output (keeping the master)
        for i in range(len(output.slides) - 1, -1, -1):
            rId = output.slides._sldIdLst[i].rId
            output.part.drop_rel(rId)
            output.slides._sldIdLst.remove(output.slides._sldIdLst[i])
        
        print(f"Copying {len(source.slides)} slides from source presentation...")
        
        # Copy slides from source to output
        for i, slide in enumerate(source.slides):
            print(f"Processing slide {i+1}...")
            
            # Choose appropriate layout
            layout_idx = 0 if i == 0 else 1  # Title slide for first, Title & Content for others
            new_slide = output.slides.add_slide(output.slide_layouts[layout_idx])
            
            # Copy text from placeholders
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                    
                # Find matching placeholder in new slide
                for placeholder in new_slide.placeholders:
                    if hasattr(placeholder, "text"):
                        try:
                            # Try to match by index or type
                            if shape.placeholder_format.idx == placeholder.placeholder_format.idx:
                                placeholder.text = shape.text
                                break
                        except (AttributeError, TypeError):
                            # If not a placeholder or other issue, skip
                            pass
                else:
                    # If we didn't find a match, put it in the first available text placeholder
                    for placeholder in new_slide.placeholders:
                        if hasattr(placeholder, "text"):
                            placeholder.text = shape.text
                            break
            
        # Save the result
        print(f"Saving presentation to: {output_path}")
        output.save(output_path)
        print("Done!")
        return True
        
    except Exception as e:
        import traceback
        print(f"ERROR: {str(e)}")
        traceback.print_exc()
        return False

def main():
    # Get the current directory
    current_dir = os.path.dirname(os.path.abspath(__file__))
    if not current_dir:
        current_dir = os.getcwd()
    
    # Paths
    source_path = os.path.join(current_dir, "AI_ML_Healthcare_Presentation.pptx")
    template_path = os.path.join(current_dir, "template", "Modern project kickoff presentation.pptx")
    output_path = os.path.join(current_dir, "AI_ML_Healthcare_Themed.pptx")
    
    print(f"Source presentation: {source_path}")
    print(f"Template presentation: {template_path}")
    print(f"Output path: {output_path}")
    
    # Apply the template
    success = copy_slides_to_template(source_path, template_path, output_path)
    
    if success:
        print(f"\nPresentation successfully themed and saved to:")
        print(output_path)
    else:
        print("\nFailed to apply template. See error messages above.")

if __name__ == "__main__":
    main()
