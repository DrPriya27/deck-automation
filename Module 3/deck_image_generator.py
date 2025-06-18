###############################################################################  
# ppt_image_generator.py  –  Generate images for PowerPoint slides  
###############################################################################  
import os, base64, mimetypes, json, textwrap, re, io, argparse, requests
from pathlib import Path
import time
import uuid
  
from dotenv import load_dotenv  
from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
  
###############################################################################  
# 1.  Azure OpenAI client  
###############################################################################  
load_dotenv()  

# Initialize Azure OpenAI client
llm = AzureChatOpenAI(
            azure_endpoint=os.getenv('AZURE_OPENAI_ENDPOINT'),
            azure_deployment="gpt-4o",
            api_version="2024-06-01",
            api_key=os.getenv('AZURE_OPENAI_API_KEY'),
            max_tokens=None,
            temperature=0.7,
)

# Use Hugging Face for image generation
import huggingface_hub
from huggingface_hub import InferenceClient

# Initialize Hugging Face client for image generation
HF_TOKEN = os.getenv("HF_TOKEN")  # Your Hugging Face API token
# Default to a stable diffusion model that's good for general purpose images
HF_MODEL_ID = os.getenv("HF_MODEL_ID", "stabilityai/stable-diffusion-xl-base-1.0")



hf_client = InferenceClient(
            provider="hf-inference",
            api_key=os.environ["HF_TOKEN"],
)
        

  
###############################################################################  
# 2.  Helper functions  
###############################################################################  
IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}  
  
def extract_slide_content(slide):
    """Extract text content from a slide"""
    content = {"title": "", "text": []}
    
    # Extract title
    if slide.shapes.title:
        content["title"] = slide.shapes.title.text
    
    # Extract text from all shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text != content["title"]:
            # Skip empty paragraphs
            if shape.text.strip():
                content["text"].append(shape.text)
    
    return content

def generate_image_prompt(slide_content):
    """Generate a prompt for image creation based on slide content"""
    title = slide_content["title"]
    text = " ".join(slide_content["text"])
    
    system_message = """You are an expert at creating detailed image prompts for DALL-E.
    Based on the PowerPoint slide content provided, create three distinct image prompts:
    1. A high-quality infographic that visualizes the key concepts
    2. A professional diagram or chart related to the topic
    3. A relevant metaphorical or conceptual illustration
    
    For each prompt, focus on creating professional, business-appropriate visuals.
    Keep each image prompt under 100 words but make it detailed and specific.
    Include style guidance like "professional", "modern", "corporate style", etc.
    DO NOT mention text that should appear in the image as DALL-E cannot reliably render text.
    """
    
    user_message = f"Create three image prompts for a slide with title: '{title}' and content: '{text}'"
    
    messages_mod = [
        SystemMessage(content=system_message),
        HumanMessage(content=user_message)
    ]

    # Get image prompt suggestions from the model
    response = llm.invoke(messages_mod)
    image_prompts = response.content.strip().split("\n\n")
    
    # Clean up and extract the actual prompts
    clean_prompts = []
    current_prompt = ""
    
    for line in response.content.strip().split("\n"):
        if line.startswith(("1.", "2.", "3.")):
            if current_prompt:
                clean_prompts.append(current_prompt.strip())
            current_prompt = line[2:].strip()
        elif line and current_prompt:
            current_prompt += " " + line.strip()
    
    if current_prompt:
        clean_prompts.append(current_prompt.strip())
    
    # Ensure we have exactly 3 prompts
    while len(clean_prompts) < 3:
        clean_prompts.append(f"Professional business infographic related to {title}")
    
    return clean_prompts[:3]

def generate_image_with_huggingface(prompt, output_path):
    """Generate an image using Hugging Face models and save it to the output path"""
    try:
        if hf_client is None:
            print("Hugging Face client not initialized. Creating placeholder image.")
            create_placeholder_image(prompt, output_path)
            return False
        
        print(f"Generating image with Hugging Face model {HF_MODEL_ID}...")
        
        # Generate the image using Hugging Face Inference API
        image = hf_client.text_to_image(
            prompt=prompt,
            model=HF_MODEL_ID,
            negative_prompt="text, watermark, signature, blurry, distorted, low quality, ugly",
            width=1024,
            height=1024,
        )
        
        
        # Save the image
        image.save(output_path)
        print(f"Image saved to {output_path}")
        return True
        
    except Exception as e:
        print(f"Error generating image with Hugging Face: {e}")
        # Create a placeholder image
        create_placeholder_image(prompt, output_path)
        return False

def create_placeholder_image(prompt, output_path):
    """Create a placeholder image with text when DALL-E is unavailable"""
    # Create a figure with a white background
    plt.figure(figsize=(10, 10), facecolor='white')
    
    # Remove axes
    plt.axis('off')
    
    # Add text to the image
    plt.text(0.5, 0.5, f"Image placeholder\n\n{prompt[:200]}...", 
             ha='center', va='center', wrap=True, fontsize=12)
    
    # Save the figure
    plt.savefig(output_path, bbox_inches='tight', pad_inches=0.5)
    plt.close()
    print(f"Placeholder image saved to {output_path}")
    return True

def process_presentation(input_pptx, output_folder):
    """Process each slide in the presentation and generate images"""
    # Create output folder if it doesn't exist
    output_path = Path(output_folder)
    output_path.mkdir(exist_ok=True, parents=True)
    
    # Load the presentation
    prs = Presentation(input_pptx)
    
    # Process each slide
    for i, slide in enumerate(prs.slides):
        print(f"\nProcessing slide {i+1} of {len(prs.slides)}")
        
        # Create folder for this slide
        slide_folder = output_path / f"slide_{i+1}"
        slide_folder.mkdir(exist_ok=True)
        
        # Extract content from the slide
        slide_content = extract_slide_content(slide)
        print(f"Slide title: {slide_content['title']}")
        
        # Generate image prompts for this slide
        image_prompts = generate_image_prompt(slide_content)
          # Generate images for each prompt
        for j, prompt in enumerate(image_prompts):
            image_path = slide_folder / f"image_{j+1}.png"
            print(f"Generating image {j+1} with prompt: {prompt[:50]}...")
            
            # Generate the image
            generate_image_with_huggingface(prompt, image_path)
            
            # Add a small delay to avoid rate limits
            time.sleep(1)
            
    print(f"\nAll slides processed. Images saved to {output_path}")
    return output_path

def main():
    # Set up argument parser
    parser = argparse.ArgumentParser(description='Generate images for PowerPoint slides')
    parser.add_argument('--input', '-i', required=True, help='Input PowerPoint file')
    parser.add_argument('--output', '-o', default='output_images', help='Output folder for images')
    parser.add_argument('--hf_token', help='Hugging Face API token (can also be set as HUGGINGFACE_TOKEN environment variable)')
    parser.add_argument('--model', default="black-forest-labs/FLUX.1-dev", 
                       help='Hugging Face model ID to use for image generation (default: stabilityai/stable-diffusion-xl-base-1.0)')
    args = parser.parse_args()
    
    # Validate input file
    input_pptx = Path(args.input)
    if not input_pptx.exists():
        print(f"Error: Input file {args.input} does not exist")
        return
    
    # Check for Hugging Face token
    global hf_client, HF_MODEL_ID
    if args.hf_token:
        os.environ["HF_TOKEN"] = args.hf_token
        hf_client = InferenceClient(token=args.hf_token)
    elif not HF_TOKEN:
        print("Warning: No Hugging Face API token provided. Will generate placeholder images.")
        print("To use Hugging Face models, set the HUGGINGFACE_TOKEN environment variable")
        print("or use the --hf_token argument.")
    
    # Update model ID if specified
    if args.model:
        HF_MODEL_ID = args.model
    
    # Process the presentation
    output_folder = process_presentation(input_pptx, args.output)
    
    print(f"Images for all slides have been generated in {output_folder}")

if __name__ == "__main__":
    main()
